from pathlib import Path

import pandas as pd
from docx import Document
from pptx import Presentation

from azure.identity import DefaultAzureCredential
from azure.ai.documentintelligence import DocumentIntelligenceClient

from utils.config import DOCUMENT_INTELLIGENCE_ENDPOINT


credential = DefaultAzureCredential()

doc_client = DocumentIntelligenceClient(
    endpoint=DOCUMENT_INTELLIGENCE_ENDPOINT,
    credential=credential,
)


def extract_with_document_intelligence(
    file_path: str,
    content_type: str,
) -> dict:
    """
    Used for:
    - PDF
    - PNG
    - JPG/JPEG
    """

    with open(file_path, "rb") as file:
        poller = doc_client.begin_analyze_document(
            model_id="prebuilt-layout",
            body=file,
            content_type=content_type,
        )

    result = poller.result()

    pages = []

    if result.pages:
        for page in result.pages:
            lines = []

            if page.lines:
                for line in page.lines:
                    lines.append(line.content)

            pages.append(
                {
                    "page_number": page.page_number,
                    "text": "\n".join(lines),
                }
            )

    return {
        "pages": pages,
        "tables": [],
        "content": result.content or "",
    }


def extract_docx(file_path: str) -> dict:
    """
    Used for Word documents (.docx)
    """

    doc = Document(file_path)

    paragraphs = [
        p.text
        for p in doc.paragraphs
        if p.text.strip()
    ]

    return {
        "pages": [],
        "tables": [],
        "content": "\n".join(paragraphs),
    }


def extract_txt(file_path: str) -> dict:
    """
    Used for text files (.txt)
    """

    with open(
        file_path,
        "r",
        encoding="utf-8",
        errors="ignore",
    ) as file:
        text = file.read()

    return {
        "pages": [],
        "tables": [],
        "content": text,
    }


def extract_csv(file_path: str) -> dict:
    """
    Used for CSV files (.csv)
    """

    df = pd.read_csv(file_path)

    return {
        "pages": [],
        "tables": [],
        "content": df.to_string(index=False),
    }


def extract_excel(file_path: str) -> dict:
    """
    Used for Excel files (.xlsx/.xls)
    """

    sheets = pd.read_excel(
        file_path,
        sheet_name=None,
    )

    content_parts = []

    for sheet_name, df in sheets.items():
        content_parts.append(f"Sheet: {sheet_name}")

        if not df.empty:
            content_parts.append(
                df.to_string(index=False)
            )

    return {
        "pages": [],
        "tables": [],
        "content": "\n\n".join(content_parts),
    }


def extract_pptx(file_path: str) -> dict:
    """
    Used for PowerPoint files (.pptx)
    """

    presentation = Presentation(file_path)

    slide_content = []

    for slide_number, slide in enumerate(
        presentation.slides,
        start=1,
    ):
        slide_content.append(
            f"Slide {slide_number}"
        )

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()

                if text:
                    slide_content.append(text)

    return {
        "pages": [],
        "tables": [],
        "content": "\n".join(slide_content),
    }


def extract_document_from_file(
    file_path: str,
) -> dict:
    """
    Automatically routes extraction
    based on file extension.
    """

    extension = Path(file_path).suffix.lower()

    # PDFs
    if extension == ".pdf":
        return extract_with_document_intelligence(
            file_path,
            "application/pdf",
        )

    # Images
    elif extension == ".png":
        return extract_with_document_intelligence(
            file_path,
            "image/png",
        )

    elif extension in [".jpg", ".jpeg"]:
        return extract_with_document_intelligence(
            file_path,
            "image/jpeg",
        )

    # Word
    elif extension == ".docx":
        return extract_docx(file_path)

    # Text
    elif extension == ".txt":
        return extract_txt(file_path)

    # CSV
    elif extension == ".csv":
        return extract_csv(file_path)

    # Excel
    elif extension in [".xlsx", ".xls"]:
        return extract_excel(file_path)

    # PowerPoint
    elif extension == ".pptx":
        return extract_pptx(file_path)

    raise ValueError(
        f"Unsupported file type: {extension}"
    )